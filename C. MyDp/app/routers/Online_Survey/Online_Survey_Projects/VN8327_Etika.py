from app.classes.AP_DataConverter import APDataConverter
from ..export_online_survey_data_table import DataTableGenerator
from ..convert_unstack import ConvertUnstack
from ..convert_stack import ConvertStack
import pandas as pd
import numpy as np
import time
import traceback
import datetime
from pptx import Presentation



class VN8327Etika(APDataConverter, DataTableGenerator, ConvertUnstack, ConvertStack):

    def convert_vn8327_etika(self, py_script_file, tables_format_file, codelist_file, coding_file):
        """DO NOT EDIT THIS FUNCTION"""

        try:
            start_time = time.time()

            if py_script_file:
                exec(py_script_file.file.read())
                exit()

            self.processing_script_vn8327_etika(tables_format_file, codelist_file, coding_file)

            self.logger.warning('DATA PROCESSING COMPLETED in %d seconds' % (time.time() - start_time))

        except Exception:
            self.logger.error(traceback.format_exc())
            str_err_log_name = self.str_file_name.replace('.xlsx', '_Errors.txt')
            with open(str_err_log_name, 'w') as err_log_txt:
                err_log_txt.writelines(traceback.format_exc())


    def processing_script_vn8327_etika(self, tables_format_file, codelist_file, coding_file):
        """
        :param tables_format_file: txt uploaded table format file
        :param codelist_file: txt uploaded codelist file
        :param coding_file: txt uploaded coding file
        :return: zip file include files: .sav, .sps, .xlsx(rawdata, topline)

        EDIT FROM HERE
        """

        # Convert system rawdata to dataframes--------------------------------------------------------------------------
        self.logger.info('Convert system rawdata to dataframes')

        self.lstDrop.extend(['Q0a_RespondentName'])

        df_data, df_info = self.convert_df_mc()

        info_col_name = ['var_name', 'var_lbl', 'var_type', 'val_lbl']

        # Pre processing------------------------------------------------------------------------------------------------
        self.logger.info('Pre processing')

        df_data.rename(columns={'S2_Age_Group': 'S2_AgeGroup'}, inplace=True)
        df_info.loc[df_info['var_name'] == 'S2_Age_Group', ['var_name']] = ['S2_AgeGroup']

        # dict_add_new_qres = {
        #     'CITY_Combined': ['CITY', 'SA', {'1': 'HCM', '2': 'HN'}, 2],
        # }
        #
        # for key, val in dict_add_new_qres.items():
        #     df_info = pd.concat([df_info, pd.DataFrame(columns=info_col_name, data=[[key, val[0], val[1], val[2]]])], axis=0, ignore_index=True)
        #     df_data = pd.concat([df_data, pd.DataFrame(columns=[key], data=[val[-1]] * df_data.shape[0])], axis=1)
        #
        # df_fil = df_data.query('Recruit_S1_Thanh_pho.isin([1, 2])')
        # df_data.loc[df_fil.index, ['CITY_Combined']] = [1]
        # del df_fil
        #
        # df_data.reset_index(drop=True, inplace=True)
        # df_info.reset_index(drop=True, inplace=True)
        #
        # # # ----------------------------------
        # # df_data.to_excel('zzz_df_data.xlsx')
        # # df_info.to_excel('zzz_df_info.xlsx')
        # # # ----------------------------------

        # --------------------------------------------------------------------------------------------------------------
        df_data_unstack, df_info_unstack = df_data.copy(), df_info.copy()
        # --------------------------------------------------------------------------------------------------------------

        dict_add_new_qres = {

            'Ma_SP_SP1': ['Mã SP', 'SA', {'101': 'CO', '102': 'CON', '103': 'YO', '104': 'PO', '105': 'CS', '106': 'CG', '107': 'PG', '108': 'PS', '109': 'GO', '110': 'GS'}, np.nan],
            'Ma_SP_SP2': ['Mã SP', 'SA', {'101': 'CO', '102': 'CON', '103': 'YO', '104': 'PO', '105': 'CS', '106': 'CG', '107': 'PG', '108': 'PS', '109': 'GO', '110': 'GS'}, np.nan],
            'Ma_SP_SP3': ['Mã SP', 'SA', {'101': 'CO', '102': 'CON', '103': 'YO', '104': 'PO', '105': 'CS', '106': 'CG', '107': 'PG', '108': 'PS', '109': 'GO', '110': 'GS'}, np.nan],

            'FC1_By_Product_Code': ['FC1_By_Product_Code', 'SA', {'101': 'CO', '102': 'CON', '103': 'YO', '104': 'PO', '105': 'CS', '106': 'CG', '107': 'PG', '108': 'PS', '109': 'GO', '110': 'GS'}, np.nan],

            'P1_101': ['P1. CO', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            'P1_102': ['P1. CON', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            'P1_103': ['P1. YO', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            'P1_104': ['P1. PO', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            'P1_105': ['P1. CS', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            'P1_106': ['P1. CG', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            'P1_107': ['P1. PG', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            'P1_108': ['P1. PS', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            'P1_109': ['P1. GO', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            'P1_110': ['P1. GS', 'SA', {'1': 'Dislike very much', '2': 'Somewhat dislike', '3': 'Neither like nor dislike it', '4': 'Somewhat like', '5': 'Like very much'}, np.nan],
            
            'FC0_YN_SP1': ['FC0. After tasting all the cultured based milk products, please select all the products that you are interested in buying', 'SA', {'1': 'Yes', '2': 'No'}, 2],
            'FC0_YN_SP2': ['FC0. After tasting all the cultured based milk products, please select all the products that you are interested in buying', 'SA', {'1': 'Yes', '2': 'No'}, 2],
            'FC0_YN_SP3': ['FC0. After tasting all the cultured based milk products, please select all the products that you are interested in buying', 'SA', {'1': 'Yes', '2': 'No'}, 2],

            'FC1_YN_SP1': ['FC1. After tasting all the cultured based milk products, which one do you intend to buy the most?', 'SA', {'1': 'Yes', '2': 'No'}, 2],
            'FC1_YN_SP2': ['FC1. After tasting all the cultured based milk products, which one do you intend to buy the most?', 'SA', {'1': 'Yes', '2': 'No'}, 2],
            'FC1_YN_SP3': ['FC1. After tasting all the cultured based milk products, which one do you intend to buy the most?', 'SA', {'1': 'Yes', '2': 'No'}, 2],

            'FC2_OE_SP1': ['FC2_OE. What are the reasons that is selected', 'FT', {}, np.nan],
            'FC2_OE_SP2': ['FC2_OE. What are the reasons that is selected', 'FT', {}, np.nan],
            'FC2_OE_SP3': ['FC2_OE. What are the reasons that is selected', 'FT', {}, np.nan],

            'FC3_OE_SP1': ['FC3_OE. What do you think are the keywords that illustrate the tastes of', 'FT', {}, np.nan],
            'FC3_OE_SP2': ['FC3_OE. What do you think are the keywords that illustrate the tastes of', 'FT', {}, np.nan],
            'FC3_OE_SP3': ['FC3_OE. What do you think are the keywords that illustrate the tastes of', 'FT', {}, np.nan],

            'FC4_OE_SP1': ['FC4_OE. What are the reasons that product is not selected', 'FT', {}, np.nan],
            'FC4_OE_SP2': ['FC4_OE. What are the reasons that product is not selected', 'FT', {}, np.nan],
            'FC4_OE_SP3': ['FC4_OE. What are the reasons that product is not selected', 'FT', {}, np.nan],

            'C3_New_1': ['C3. Original', 'SA', {'1': 'B2B', '2': 'Medium', '3': 'T2B'}, np.nan],
            'C3_New_2': ['C3. Grape', 'SA', {'1': 'B2B', '2': 'Medium', '3': 'T2B'}, np.nan],
            'C3_New_3': ['C3. Strawberry', 'SA', {'1': 'B2B', '2': 'Medium', '3': 'T2B'}, np.nan],

            'PI5_New_1': ['PI5. Original', 'SA', {'1': 'B2B', '2': 'Medium', '3': 'T2B'}, np.nan],
            'PI5_New_2': ['PI5. Original Nat de Coco', 'SA', {'1': 'B2B', '2': 'Medium', '3': 'T2B'}, np.nan],
            'PI5_New_3': ['PI5. Calpis_Grape', 'SA', {'1': 'B2B', '2': 'Medium', '3': 'T2B'}, np.nan],
            'PI5_New_4': ['PI5. Calpis_Strawberry', 'SA', {'1': 'B2B', '2': 'Medium', '3': 'T2B'}, np.nan],
            'PI5_New_5': ['PI5. Goodday_Original', 'SA', {'1': 'B2B', '2': 'Medium', '3': 'T2B'}, np.nan],
            'PI5_New_6': ['PI5. Goodday_Strawberry', 'SA', {'1': 'B2B', '2': 'Medium', '3': 'T2B'}, np.nan],

            'P2_New_1': ['P2. Assuming this drink was available at a reasonable price where you normally buy, how likely would you be to buy it?', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'P2_New_2': ['P2. Assuming this drink was available at a reasonable price where you normally buy, how likely would you be to buy it?', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'P2_New_3': ['P2. Assuming this drink was available at a reasonable price where you normally buy, how likely would you be to buy it?', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'P2_New_4': ['P2. Assuming this drink was available at a reasonable price where you normally buy, how likely would you be to buy it?', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'P2_New_5': ['P2. Assuming this drink was available at a reasonable price where you normally buy, how likely would you be to buy it?', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'P2_New_6': ['P2. Assuming this drink was available at a reasonable price where you normally buy, how likely would you be to buy it?', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],

            'PI2_New_1': ['PI2. Original', 'SA', {'1': '18,000 VND', '2': '19,000 VND', '3': '20,000 VND', '4': '21,000 VND', '5': '22,000 VND', '6': '23, 000 VND', '7': '24,000 VND', '8': '25,000 VND', '9': '26,000 VND', '10': '27,000 VND', '11': '28,000 VND', '12': '29,000 VND', '13': '30,000 VND'}, np.nan],
            'PI2_New_2': ['PI2. Original Nat de Coco', 'SA', {'1': '18,000 VND', '2': '19,000 VND', '3': '20,000 VND', '4': '21,000 VND', '5': '22,000 VND', '6': '23, 000 VND', '7': '24,000 VND', '8': '25,000 VND', '9': '26,000 VND', '10': '27,000 VND', '11': '28,000 VND', '12': '29,000 VND', '13': '30,000 VND'}, np.nan],
            'PI2_New_3': ['PI2. Calpis_Grape', 'SA', {'1': '18,000 VND', '2': '19,000 VND', '3': '20,000 VND', '4': '21,000 VND', '5': '22,000 VND', '6': '23, 000 VND', '7': '24,000 VND', '8': '25,000 VND', '9': '26,000 VND', '10': '27,000 VND', '11': '28,000 VND', '12': '29,000 VND', '13': '30,000 VND'}, np.nan],
            'PI2_New_4': ['PI2. Calpis_Strawberry', 'SA', {'1': '18,000 VND', '2': '19,000 VND', '3': '20,000 VND', '4': '21,000 VND', '5': '22,000 VND', '6': '23, 000 VND', '7': '24,000 VND', '8': '25,000 VND', '9': '26,000 VND', '10': '27,000 VND', '11': '28,000 VND', '12': '29,000 VND', '13': '30,000 VND'}, np.nan],
            'PI2_New_5': ['PI2. Goodday_Original', 'SA', {'1': '18,000 VND', '2': '19,000 VND', '3': '20,000 VND', '4': '21,000 VND', '5': '22,000 VND', '6': '23, 000 VND', '7': '24,000 VND', '8': '25,000 VND', '9': '26,000 VND', '10': '27,000 VND', '11': '28,000 VND', '12': '29,000 VND', '13': '30,000 VND'}, np.nan],
            'PI2_New_6': ['PI2. Goodday_Strawberry', 'SA', {'1': '18,000 VND', '2': '19,000 VND', '3': '20,000 VND', '4': '21,000 VND', '5': '22,000 VND', '6': '23, 000 VND', '7': '24,000 VND', '8': '25,000 VND', '9': '26,000 VND', '10': '27,000 VND', '11': '28,000 VND', '12': '29,000 VND', '13': '30,000 VND'}, np.nan],

            'PI5_New2_1': ['PI5. Original', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'PI5_New2_2': ['PI5. Original Nat de Coco', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'PI5_New2_3': ['PI5. Calpis_Grape', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'PI5_New2_4': ['PI5. Calpis_Strawberry', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'PI5_New2_5': ['PI5. Goodday_Original', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'PI5_New2_6': ['PI5. Goodday_Strawberry', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],

            'P2_Price1': ['P2. Assuming this drink was available at a reasonable price where you normally buy, how likely would you be to buy it?', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],
            'P2_Price2': ['P2. Assuming this drink was available at a reasonable price where you normally buy, how likely would you be to buy it?', 'SA', {'1': 'Definitely would not buy', '2': 'Probably would not buy', '3': 'Might or might not buy', '4': 'Probably would buy', '5': 'Definitely would buy'}, np.nan],

            'P2_Compare_PI5': ['P2 vs PI5. Compare purchase intention', 'SA', {'1': 'Different both products', '2': 'Different for one products', '3': 'Same for both products'}, 0],

            'P2_Compare_PI2_PI5': ['P2 vs PI5 with PI2. Users who change their purchase intention after hearing the pricing appropriately', 'SA', {
                '1': 'Purchase intention with pricing > purchase intention without pricing, penetration price > 20,000 VND',
                '2': 'Purchase intention with pricing > purchase intention without pricing, penetration price <=20,000 VND',
                '99': 'Else',
            }, np.nan],

            'P2_Compare_PI2_PI5_v2': ['P2 vs PI5 with PI2. Users who change their purchase intention after hearing the pricing appropriately: Purchase intention with pricing > purchase intention without pricing, penetration price <=20,000 VND', 'SA', {
                '1': 'Change 1 product',
                '2': 'Change 2 products',
                '99': 'Else',
            }, 0],

            'P2_Compare_PI5_Price1': ['P2 vs PI5. Compare purchase intention', 'SA', {'1': 'Change their purchase intention after hearing the pricing appropriately', '2': 'Not change their purchase intention after hearing the pricing appropriately'}, np.nan],
            'P2_Compare_PI5_Price2': ['P2 vs PI5. Compare purchase intention', 'SA', {'1': 'Change their purchase intention after hearing the pricing appropriately', '2': 'Not change their purchase intention after hearing the pricing appropriately'}, np.nan],

            'P2_Compare_PI5_2_Price1': ['P2 vs PI5. Compare purchase intention v2', 'SA', {'1': 'Not buy in P2(choose 1/2/3) change to buy in PI5(choose 4/5)', '2': 'Buy in P2(choose 4/5) change to not buy in PI5(choose 1/2/3)', '3': 'No change'}, np.nan],
            'P2_Compare_PI5_2_Price2': ['P2 vs PI5. Compare purchase intention v2', 'SA', {'1': 'Not buy in P2(choose 1/2/3) change to buy in PI5(choose 4/5)', '2': 'Buy in P2(choose 4/5) change to not buy in PI5(choose 1/2/3)', '3': 'No change'}, np.nan],




        }

        for key, val in dict_add_new_qres.items():
            df_info = pd.concat([df_info, pd.DataFrame(columns=info_col_name, data=[[key, val[0], val[1], val[2]]])], axis=0, ignore_index=True)
            df_data = pd.concat([df_data, pd.DataFrame(columns=[key], data=[val[-1]] * df_data.shape[0])], axis=1)

        for i in range(1, 4):
            df_data[f'Ma_SP_PA3_1_SP{i}'].replace({1: 101, 2: 102, 3: 103, 4: 104}, inplace=True)
            df_data[f'Ma_SP_PA3_2_SP{i}'].replace({1: 105, 2: 106, 3: 107, 4: 108}, inplace=True)

            if i <= 2:
                df_data[f'Ma_SP_PA3_3_SP{i}'].replace({1: 109, 2: 103, 3: 104}, inplace=True)
                df_data[f'Ma_SP_PA3_4_SP{i}'].replace({1: 110, 2: 108, 3: 107}, inplace=True)

                df_data[f'Ma_SP_SP{i}'] = df_data[[f'Ma_SP_PA3_1_SP{i}', f'Ma_SP_PA3_2_SP{i}', f'Ma_SP_PA3_3_SP{i}', f'Ma_SP_PA3_4_SP{i}']].sum(axis=1)

            else:
                df_data[f'Ma_SP_SP{i}'] = df_data[[f'Ma_SP_PA3_1_SP{i}', f'Ma_SP_PA3_2_SP{i}']].sum(axis=1)

            df_fil = df_data.query(f"FC0_2Products_1 == {i} | FC0_2Products_2 == {i} | FC0_3Products_1 == {i} | FC0_3Products_2 == {i} | FC0_3Products_3 == {i}")
            df_data.loc[df_fil.index, [f'FC0_YN_SP{i}']] = [1]
            del df_fil

            df_fil = df_data.query(f"FC1_2Products == {i} | FC1_3Products == {i}")
            df_data.loc[df_fil.index, [f'FC1_YN_SP{i}']] = [1]
            del df_fil

            df_fil = df_data.query(f"FC1_YN_SP{i} == 1")
            df_data.loc[df_fil.index, f'FC2_OE_SP{i}'] = df_data.loc[df_fil.index, 'FC2_OE']
            df_data.loc[df_fil.index, f'FC3_OE_SP{i}'] = df_data.loc[df_fil.index, 'FC3_OE']
            del df_fil

        df_fil = df_data.query("FC1_YN_SP1 == 1")
        df_data.loc[df_fil.index, 'FC4_OE_SP2'] = df_data.loc[df_fil.index, 'FC4_OE_o1']
        df_data.loc[df_fil.index, 'FC4_OE_SP3'] = df_data.loc[df_fil.index, 'FC4_OE_o2']

        df_data.loc[df_fil.index, 'FC1_By_Product_Code'] = df_data.loc[df_fil.index, 'Ma_SP_SP1']

        del df_fil

        df_fil = df_data.query("FC1_YN_SP2 == 1")
        df_data.loc[df_fil.index, 'FC4_OE_SP1'] = df_data.loc[df_fil.index, 'FC4_OE_o1']
        df_data.loc[df_fil.index, 'FC4_OE_SP3'] = df_data.loc[df_fil.index, 'FC4_OE_o2']

        df_data.loc[df_fil.index, 'FC1_By_Product_Code'] = df_data.loc[df_fil.index, 'Ma_SP_SP2']

        del df_fil

        df_fil = df_data.query("FC1_YN_SP3 == 1")
        df_data.loc[df_fil.index, 'FC4_OE_SP1'] = df_data.loc[df_fil.index, 'FC4_OE_o1']
        df_data.loc[df_fil.index, 'FC4_OE_SP2'] = df_data.loc[df_fil.index, 'FC4_OE_o2']

        df_data.loc[df_fil.index, 'FC1_By_Product_Code'] = df_data.loc[df_fil.index, 'Ma_SP_SP3']

        del df_fil

        for i in range(1, 4):
            for j in range(1, 3):
                df_fil = df_data.query(f"Ma_Pack{j} == {i}")
                df_data.loc[df_fil.index, f'C3_New_{i}'] = df_data.loc[df_fil.index, f'C3_Pack{j}']
                del df_fil

        for i in range(1, 4):
            df_data[f'C3_New_{i}'].replace({1: 1, 2: 1, 3: 2, 4: 3, 5: 3}, inplace=True)

        for i in range(1, 7):
            for j in range(1, 3):
                df_fil = df_data.query(f"Ma_Price{j} == {i}")
                df_data.loc[df_fil.index, f'PI5_New_{i}'] = df_data.loc[df_fil.index, f'PI5_Price{j}']
                df_data.loc[df_fil.index, f'PI5_New2_{i}'] = df_data.loc[df_fil.index, f'PI5_Price{j}']
                df_data.loc[df_fil.index, f'PI2_New_{i}'] = df_data.loc[df_fil.index, f'PI2_Price{j}']
                del df_fil

        for i in range(1, 7):
            df_data[f'PI5_New_{i}'].replace({1: 1, 2: 1, 3: 2, 4: 3, 5: 3}, inplace=True)


        dict_compare = {101: 1, 102: 2, 106: 3, 105: 4, 109: 5, 110: 6}

        for i in range(1, 4):
            for sp in [101, 102, 103, 104, 105, 106, 107, 108, 109, 110]:
                df_fil = df_data.query(f"Ma_SP_SP{i} == {sp}")
                df_data.loc[df_fil.index, f'P1_{sp}'] = df_data.loc[df_fil.index, f'P1_SP{i}']
                del df_fil

                if sp in [101, 102, 105, 106, 109, 110]:
                    df_fil = df_data.query(f"Ma_SP_SP{i} == {sp}")
                    df_data.loc[df_fil.index, f'P2_New_{dict_compare.get(sp)}'] = df_data.loc[df_fil.index, f'P2_SP{i}']
                    del df_fil


        for i in [1, 2]:
            for price_code in range(1, 7):
                df_fil = df_data.query(f"Ma_Price{i} == {price_code}")
                df_data.loc[df_fil.index, f'P2_Price{i}'] = df_data.loc[df_fil.index, f'P2_New_{price_code}']
                del df_fil



        # 'P2_Compare_PI5': ['P2 vs PI5. Compare purchase intention', 'SA',
        # {'1': 'Different both products',
        # '2': 'Different for one products',
        # '3': 'Same for both products'}, np.nan],
        for i in range(1, 7):
            df_fil = df_data.query(f"(PI5_New2_{i} != P2_New_{i}) & ~PI5_New2_{i}.isnull() & ~P2_New_{i}.isnull()")
            df_data.loc[df_fil.index, 'P2_Compare_PI5'] = df_data.loc[df_fil.index, 'P2_Compare_PI5'] + 1
            del df_fil

        df_data.loc[df_data.eval("P2_Compare_PI5 == 0"), 'P2_Compare_PI5'] = 3


        # 'P2_Compare_PI2_PI5':
        # ['P2 vs PI5 with PI2. Compare purchase intention', 'SA',
        # {'1': 'users who change their purchase intention after hearing the pricing appropriately (purchase intention with pricing > purchase intention without pricing, penetration price > 20,000 VND)',
        # '2': 'users who change their purchase intention after hearing the pricing inappropriately (purchase intention with pricing > purchase intention without pricing, penetration price <=20,000 VND)'}, np.nan],

        for i in range(1, 7):
            df_fil = df_data.query(f"(PI5_New2_{i} > P2_New_{i} & PI2_New_{i}.isin([4, 5, 6, 7, 8, 9, 10, 11, 12, 13])) & ~PI2_New_{i}.isnull() & ~PI5_New2_{i}.isnull() & ~P2_New_{i}.isnull()")
            df_data.loc[df_fil.index, 'P2_Compare_PI2_PI5'] = 1
            del df_fil

        for i in range(1, 7):
            df_fil = df_data.query(f"(PI5_New2_{i} > P2_New_{i} & PI2_New_{i}.isin([1, 2, 3])) & ~PI2_New_{i}.isnull() & ~PI5_New2_{i}.isnull() & ~P2_New_{i}.isnull()")
            df_data.loc[df_fil.index, 'P2_Compare_PI2_PI5'] = 2
            del df_fil

        df_data.loc[df_data.eval("P2_Compare_PI2_PI5.isnull()"), 'P2_Compare_PI2_PI5'] = 99

        # 'P2_Compare_PI2_PI5_v2': [
        #     'P2 vs PI5 with PI2. Users who change their purchase intention after hearing the pricing appropriately: Purchase intention with pricing > purchase intention without pricing, penetration price <=20,000 VND',
        #     'SA', {
        #         '1': 'Change 1 product',
        #         '2': 'Change 2 products',
        #         '99': 'Not Change',
        #     }, 0],
        for i in range(1, 7):
            df_fil = df_data.query(f"(PI5_New2_{i} > P2_New_{i} & PI2_New_{i}.isin([1, 2, 3])) & ~PI2_New_{i}.isnull() & ~PI5_New2_{i}.isnull() & ~P2_New_{i}.isnull()")
            df_data.loc[df_fil.index, 'P2_Compare_PI2_PI5_v2'] = df_data.loc[df_fil.index, 'P2_Compare_PI2_PI5_v2'] + 1
            del df_fil

        df_data.loc[df_data.eval("P2_Compare_PI2_PI5_v2 == 0"), 'P2_Compare_PI2_PI5_v2'] = 99




        # P2_Compare_PI5_Price1 | P2_Compare_PI5_Price2
        # {'1': 'Change their purchase intention after hearing the pricing appropriately',
        # '2': 'Not change their purchase intention after hearing the pricing appropriately'}

        # 'P2_Compare_PI5_2_Price1': ['P2 vs PI5. Compare purchase intention v2', 'SA', {'1': 'Change not buy in P2 to buy in P5', '2': 'Change buy in P2 to not buy in P5', '3': 'No change'}, np.nan],
        # 'P2_Compare_PI5_2_Price2': ['P2 vs PI5. Compare purchase intention v2', 'SA', {'1': 'Change not buy in P2 to buy in P5', '2': 'Change buy in P2 to not buy in P5', '3': 'No change'}, np.nan],

        for i in [1, 2]:
            df_fil = df_data.query(f"(P2_Price{i} != PI5_Price{i}) & ~P2_Price{i}.isnull() & ~PI5_Price{i}.isnull()")
            df_data.loc[df_fil.index, f'P2_Compare_PI5_Price{i}'] = 1
            del df_fil

            df_fil = df_data.query(f"(P2_Price{i} == PI5_Price{i}) & ~P2_Price{i}.isnull() & ~PI5_Price{i}.isnull()")
            df_data.loc[df_fil.index, f'P2_Compare_PI5_Price{i}'] = 2
            del df_fil

            df_fil = df_data.query(f"(P2_Price{i}.isin([1, 2, 3]) & PI5_Price{i}.isin([4, 5])) & ~P2_Price{i}.isnull() & ~PI5_Price{i}.isnull()")
            df_data.loc[df_fil.index, f'P2_Compare_PI5_2_Price{i}'] = 1
            del df_fil

            df_fil = df_data.query(f"(P2_Price{i}.isin([4, 5]) & PI5_Price{i}.isin([1, 2, 3])) & ~P2_Price{i}.isnull() & ~PI5_Price{i}.isnull()")
            df_data.loc[df_fil.index, f'P2_Compare_PI5_2_Price{i}'] = 2
            del df_fil

            df_fil = df_data.query(f"((P2_Price{i}.isin([1, 2, 3]) & PI5_Price{i}.isin([1, 2, 3])) | (P2_Price{i}.isin([4, 5]) & PI5_Price{i}.isin([4, 5]))) & ~P2_Price{i}.isnull() & ~PI5_Price{i}.isnull()")
            df_data.loc[df_fil.index, f'P2_Compare_PI5_2_Price{i}'] = 3
            del df_fil


        df_data.reset_index(drop=True, inplace=True)
        df_info.reset_index(drop=True, inplace=True)


        # # ----------------------------------
        # df_data.to_excel('zzz_df_data.xlsx')
        # df_info.to_excel('zzz_df_info.xlsx')
        # # ----------------------------------

        # Define structure----------------------------------------------------------------------------------------------
        self.logger.info('Define structure')

        id_col = 'Q0a_RespondentID'

        sp_col = 'Ma_SP'
        pack_col = 'Ma_Pack'
        price_col = 'Ma_Price'

        lst_scr = [
            'S1',
            'S2',
            'S2_AgeGroup',
            'S3_P6M',
            'S3_Banned_1',
            'S3_Banned_2',
            'S3_Banned_3',
            'S3_Banned_4',
            'S3_Banned_5',
            'S3_Banned_6',
            'S3_marital_stt',
            'S4',
            'S4_o7',
            'S5',
            'S5_o3',
            'S6',
            'S7',
            'S8_1',
            'S8_2',
            'S8_3',
            'S8_4',
            'S8_5',
            'S8_6',
            'S8_7',
            'S8_8',
            'S8_9',
            'S8_10',
            'S8_11',
            'S8_o11',
            'S9',
            'S10_1',
            'S10_2',
            'S10_3',
            'S10_4',
            'S10_5',
            'S10_6',
            'S10_7',
            'S10_8',
            'S10_9',
            'S10_10',
            'S10_11',
            'S10_12',
            'S10_13',
            'S10_o13',
            'S11',
            'S12_1',
            'S12_2',
            'S12_3',
            'S12_4',
            'S12_5',
            'S12_6',
            'S12_7',
            'S12_8',
            'S12_9',
            'S12_10',
            'S12_11',
            'S12_12',
            'S12_13',
            'S12_14',
            'S13',
            'S14_1',
            'S14_2',
            'S14_3',
            'S14_4',
            'S14_5',
            'S14_6',
            'S14_7',
            'S14_8',
            'S14_9',
            'S15',
            'B1_1',
            'B1_2',
            'B1_3',
            'B1_4',
            'B1_5',
            'B1_6',
            'B1_7',
            'B1_8',
            'B1_9',
            'B1_o9',
            'B2_1',
            'B2_2',
            'B2_3',
            'B2_4',
            'B2_5',
            'B2_6',
            'B2_7',
            'B2_8',
            'B2_9',
            'B2_o9',
            'B3_1',
            'B3_2',
            'B3_3',
            'B3_4',
            'B3_5',
            'B3_o5',
            'B5_1',
            'B5_2',
            'B5_3',
            'B5_4',
            'B5_5',
            'B5_o5',
            'B6_1',
            'B6_2',
            'B6_3',
            'B6_4',
            'B6_o4',
            'B7_1',
            'B7_2',
            'B7_3',
            'B7_4',
            'B7_5',
            'B7_o5',
            'CO1',
            'CO2',
            'CO3',
            'ROTATION',
            'PA1',
            'PA2',
            'PA3_1',
            'PA3_2',
            'PA3_3',
            'PA3_4',

            'C3_New_1',
            'C3_New_2',
            'C3_New_3',

            'PI5_New_1',
            'PI5_New_2',
            'PI5_New_3',
            'PI5_New_4',
            'PI5_New_5',
            'PI5_New_6',

            'FC1_By_Product_Code',
            'P1_101',
            'P1_102',
            'P1_103',
            'P1_104',
            'P1_105',
            'P1_106',
            'P1_107',
            'P1_108',
            'P1_109',
            'P1_110',

            'P2_New_1',
            'P2_New_2',
            'P2_New_3',
            'P2_New_4',
            'P2_New_5',
            'P2_New_6',

            'PI2_New_1',
            'PI2_New_2',
            'PI2_New_3',
            'PI2_New_4',
            'PI2_New_5',
            'PI2_New_6',

            'PI5_New2_1',
            'PI5_New2_2',
            'PI5_New2_3',
            'PI5_New2_4',
            'PI5_New2_5',
            'PI5_New2_6',

            'P2_Compare_PI5',
            'P2_Compare_PI2_PI5',
            'P2_Compare_PI2_PI5_v2',
        ]

        dict_main = {
            1: {
                'Ma_SP_SP1': 'Ma_SP',
                'P1_SP1': 'P1',
                'P1_1_OE_SP1': 'P1_1_OE',
                'P1_2_OE_SP1': 'P1_2_OE',
                'P2_SP1': 'P2',
                'P3_SP1': 'P3',
                'P4_SP1': 'P4',
                'P4_JAR_SP1': 'P4_JAR',
                'P5_SP1': 'P5',
                'P5_JAR_SP1': 'P5_JAR',
                'P6_SP1': 'P6',
                'P6_JAR_SP1': 'P6_JAR',
                'P7_SP1': 'P7',
                'P7_JAR_SP1': 'P7_JAR',
                'P8_SP1': 'P8',
                'P8_JAR_SP1': 'P8_JAR',
                'P9_SP1': 'P9',
                'P10_SP1': 'P10',
                'FC0_YN_SP1': 'FC0_YN',
                'FC1_YN_SP1': 'FC1_YN',
                'FC2_OE_SP1': 'FC2_OE_New',
                'FC3_OE_SP1': 'FC3_OE_New',
                'FC4_OE_SP1': 'FC4_OE_New',
            },
            2: {
                'Ma_SP_SP2': 'Ma_SP',
                'P1_SP2': 'P1',
                'P1_1_OE_SP2': 'P1_1_OE',
                'P1_2_OE_SP2': 'P1_2_OE',
                'P2_SP2': 'P2',
                'P3_SP2': 'P3',
                'P4_SP2': 'P4',
                'P4_JAR_SP2': 'P4_JAR',
                'P5_SP2': 'P5',
                'P5_JAR_SP2': 'P5_JAR',
                'P6_SP2': 'P6',
                'P6_JAR_SP2': 'P6_JAR',
                'P7_SP2': 'P7',
                'P7_JAR_SP2': 'P7_JAR',
                'P8_SP2': 'P8',
                'P8_JAR_SP2': 'P8_JAR',
                'P9_SP2': 'P9',
                'P10_SP2': 'P10',
                'FC0_YN_SP2': 'FC0_YN',
                'FC1_YN_SP2': 'FC1_YN',
                'FC2_OE_SP2': 'FC2_OE_New',
                'FC3_OE_SP2': 'FC3_OE_New',
                'FC4_OE_SP2': 'FC4_OE_New',
            },
            3: {
                'Ma_SP_SP3': 'Ma_SP',
                'P1_SP3': 'P1',
                'P1_1_OE_SP3': 'P1_1_OE',
                'P1_2_OE_SP3': 'P1_2_OE',
                'P2_SP3': 'P2',
                'P3_SP3': 'P3',
                'P4_SP3': 'P4',
                'P4_JAR_SP3': 'P4_JAR',
                'P5_SP3': 'P5',
                'P5_JAR_SP3': 'P5_JAR',
                'P6_SP3': 'P6',
                'P6_JAR_SP3': 'P6_JAR',
                'P7_SP3': 'P7',
                'P7_JAR_SP3': 'P7_JAR',
                'P8_SP3': 'P8',
                'P8_JAR_SP3': 'P8_JAR',
                'P9_SP3': 'P9',
                'P10_SP3': 'P10',
                'FC0_YN_SP3': 'FC0_YN',
                'FC1_YN_SP3': 'FC1_YN',
                'FC2_OE_SP3': 'FC2_OE_New',
                'FC3_OE_SP3': 'FC3_OE_New',
                'FC4_OE_SP3': 'FC4_OE_New',
            },
        }

        lst_fc = []

        dict_pack = {
            1: {
                'Ma_Pack1': 'Ma_Pack',
                'C1_OE_Pack1': 'C1_OE',
                'C2_OE_Pack1': 'C2_OE',
                'C3_Pack1': 'C3',
                'C4_Pack1': 'C4',
                'C5_Pack1': 'C5',
                'C6_Pack1_1': 'C6_1',
                'C6_Pack1_2': 'C6_2',
                'C6_Pack1_3': 'C6_3',
                'C6_Pack1_4': 'C6_4',
                'C6_Pack1_5': 'C6_5',
                'C6_Pack1_6': 'C6_6',
                'C6_Pack1_7': 'C6_7',
                'C6_Pack1_8': 'C6_8',
                'C6_Pack1_9': 'C6_9',
                'C6_Pack1_o9': 'C6_o9',
                'C7_Pack1': 'C7',
                'C8_Pack1': 'C8',
                'C10_Pack1': 'C10',
                'C10_1_OE_Pack1': 'C10_1_OE',
                'C10_2_OE_Pack1': 'C10_2_OE',
                'C11_OE_Pack1': 'C11_OE',
            },
            2: {
                'Ma_Pack2': 'Ma_Pack',
                'C1_OE_Pack2': 'C1_OE',
                'C2_OE_Pack2': 'C2_OE',
                'C3_Pack2': 'C3',
                'C4_Pack2': 'C4',
                'C5_Pack2': 'C5',
                'C6_Pack2_1': 'C6_1',
                'C6_Pack2_2': 'C6_2',
                'C6_Pack2_3': 'C6_3',
                'C6_Pack2_4': 'C6_4',
                'C6_Pack2_5': 'C6_5',
                'C6_Pack2_6': 'C6_6',
                'C6_Pack2_7': 'C6_7',
                'C6_Pack2_8': 'C6_8',
                'C6_Pack2_9': 'C6_9',
                'C6_Pack2_o9': 'C6_o9',
                'C7_Pack2': 'C7',
                'C8_Pack2': 'C8',
                'C10_Pack2': 'C10',
                'C10_1_OE_Pack2': 'C10_1_OE',
                'C10_2_OE_Pack2': 'C10_2_OE',
                'C11_OE_Pack2': 'C11_OE',
            },
        }

        dict_price = {
            1: {
                'Ma_Price1': 'Ma_Price',
                'P2_Price1': 'P2_Price',
                'PI1_Price1': 'PI1',
                'PI2_Price1': 'PI2',
                'PI3_Price1': 'PI3',
                'PI4_Price1': 'PI4',
                'PI5_Price1': 'PI5',
                'PI6_OE_Price1': 'PI6_OE',
                'P2_Compare_PI5_Price1': 'P2_Compare_PI5_Price',
                'P2_Compare_PI5_2_Price1': 'P2_Compare_PI5_2_Price',
            },
            2: {
                'Ma_Price2': 'Ma_Price',
                'P2_Price2': 'P2_Price',
                'PI1_Price2': 'PI1',
                'PI2_Price2': 'PI2',
                'PI3_Price2': 'PI3',
                'PI4_Price2': 'PI4',
                'PI5_Price2': 'PI5',
                'PI6_OE_Price2': 'PI6_OE',
                'P2_Compare_PI5_Price2': 'P2_Compare_PI5_Price',
                'P2_Compare_PI5_2_Price2': 'P2_Compare_PI5_2_Price',
            },
        }

        # -----------------
        dict_qre_group_mean = {
            # 'Main_Q1_OL': {
            #     'range': [],  # f'0{i}' for i in range(1, 4)
            #     'group': {'cats': {1: 'B2B', 2: 'Medium', 3: 'T2B'}, 'recode': {1: 1, 2: 1, 3: 2, 4: 3, 5: 3}},
            #     'mean': {1: 1, 2: 2, 3: 3, 4: 4, 5: 5}
            # },
        }

        for qre in [
            'CO1',
            'CO2',
            'CO3',
            'P1',
            'P2',
            'P4',
            'P4_JAR',
            'P5',
            'P5_JAR',
            'P6',
            'P6_JAR',
            'P7',
            'P7_JAR',
            'P8',
            'P8_JAR',
            'P9',
            'P10',
            'C3',
            'C4',
            'C5',
            'P2_Price',
            'PI5',
        ]:
            dict_qre_group_mean[qre] = {
                'range': [],
                'group': {'cats': {1: 'B2B', 2: 'Medium', 3: 'T2B'}, 'recode': {1: 1, 2: 1, 3: 2, 4: 3, 5: 3}},
                'mean': {1: 1, 2: 2, 3: 3, 4: 4, 5: 5}
            }


        # dict_qre_OE_info_org = {'Main_Q3_Cai_thien_New_OE|1-3': ['Q3. lbl', 'MA',  {'net_code': {'999': 'Không có/Không cần', '90001|SỢI PHỞ (NET)': {'101': 'Thêm nội dung về sợi phở/làm nổi bật/nhấn mạnh ý tưởng về sợi phở', '102': 'Được làm từ gạo nếp nương tạo cảm giác thiên nhiên/tự nhiên/sạch/không thuốc', '103': 'Được làm từ gạo nếp nương giúp sợi phở thanh lành (hơn)', '104': 'Được làm từ gạo nếp nương vùng tây bắc ', '105': 'Sợi phở dẻo bùi', '106': 'Sợi phở thơm', '107': 'Sợi phở ngon', '108': 'Đề cập đến độ dai', '109': 'Gạo nếp nương đặc sắc của vùng cao tạo cảm giác thoải mái', '110': 'Gạo nếp nương đặc sắc của vùng cao tạo cảm giác thanh nhẹ', '111': 'Gạo nếp nương đặc sắc của vùng cao tạo cảm giác món phở có vị đậm đà'}, '90002|NGUYÊN LIỆU (NET)': {'201': 'Đề cập chi tiết hơn về các nguyên liệu truyền thống', '202': 'Đề cập về thịt', '203': 'Đông trùng hạ thảo giúp cân bằng cảm xúc'}, '90003|NƯỚC DÙNG (NET)': {'301': 'Được nấu theo công thức của quán phở Thìn Bờ Hồ danh tiếng 70 năm làm cho vị sợi phở ngon hấp dẫn', '302': 'Được nấu từ công thức đặc biệt của quán phở trứ danh của Hà Nội tạo cảm giác hương vị ngon, sợi phở thanh lành', '303': 'Sẽ có/mang lại vị thực tế hơn', '304': 'Được nấu theo công thức quản phở Bờ Hồ lâu đời', '305': 'Thêm nội dung về nước dùng/làm nổi bật/nhấn mạnh ý tưởng về nước dùng', '306': 'Được nấu theo công thức quản phở nổi tiếng'}, '90004|SỢI PHỞ KẾT HỢP NƯỚC DÙNG (NET)': {'401': 'Tạo sự kết hợp ấn tượng giữa sợi phở và nước dùng'}, '90005|KHÁC (NET)': {'601': 'Bổ sung thêm lợi ích/chất dinh dưỡng mang lại cho cơ thể', '610': 'Bổ sung thêm công thức nước phở Thìn đã có ngay trong gói phở ăn liền giúp tiết kiệm thời gian hơn khi ăn ngoài hàng'}, '90006|ĐIỀU CHỈNH NỘI DUNG (NET)': {'602': 'Kết hợp độc đáo đổi thành kết hợp hài hòa sẽ hay hơn', '603': 'Gạo nếp nương của vùng cao đổi thành gạo nếp nương đặc sản của vùng cao', '604': 'Giúp cân bằng cảm xúc ngay cả lúc mệt mỏi, cho cơ thể nhẹ nhàng và tạo nên một lối sống lành mạnh, thư thái đổi thành sản phẩm mang đến sức khỏe tốt cho cơ thể, giúp cơ thể hồi phục sức nhanh'}}}],}
        dict_qre_OE_info_org = dict()
        if codelist_file.filename:
            dict_qre_OE_info_org = eval(codelist_file.file.read())


        dict_qre_OE_info = dict()
        for k, v in dict_qre_OE_info_org.items():
            oe_name, oe_num_col = k.rsplit('|', 1)
            oe_num_col = oe_num_col.split('-')
            oe_num_col = range(int(oe_num_col[0]), int(oe_num_col[1]) + 1)

            for i in oe_num_col:
                dict_qre_OE_info.update({
                    f'{oe_name}_{i}': [f'{oe_name}_{i}'] + v
                })


        # lst_addin_OE_value = [['Q0a_RespondentID', 1056, 'Main_Ma_san_pham_Q3', 1, 'Main_Q5a_OE_Ly_do_thich_Y1_1', 101],]
        lst_addin_OE_value = list()
        if coding_file.filename:
            lst_addin_OE_value = eval(coding_file.file.read())


        # net code for exist qres
        dict_qre_net_info = {
            # 'Q1a_[0-9]+': {
            #     'net_code': {
            #         '900001|net|GIA VỊ': {
            #             '1': 'Tương ớt',
            #             '2': 'Nước mắm',
            #             '3': 'Nước tương',
            #             '4': 'Tương cà',
            #             '5': 'Hạt nêm',
            #         },
            #         '900002|net|DẦU GỘI/SỮA TẮM/BỘT GIẶT': {
            #             '6': 'Dầu gội',
            #             '7': 'Sữa tắm',
            #             '8': 'Bột giặt',
            #         },
            #         '900003|net|THỊT ĐÃ CHẾ BIẾN': {
            #             '9': 'Thịt hộp',
            #             '10': 'Xúc xích thanh trùng',
            #         },
            #         '900004|net|SỮA HẠT SÔ-CÔ-LA / SỮA SÔ-CÔ-LA CHO BÉ': {
            #             '11': 'Sữa hạt sô-cô-la / sữa sô-cô-la cho con',
            #         },
            #         '900005|net|THỰC PHẨM TIỆN LỢI': {
            #             '12': 'Mì ăn liền',
            #             '13': 'Phở ăn liền',
            #         },
            #     }
            # },
        }

        lst_clear_ma_oe_value = {
            # "Main_Q2c_SA_Impressive_Concept_A": [14],
            # "Main_Q2c_SA_Impressive_Concept_B": [13],
            # "Main_Q2c_SA_Impressive_Concept_C": [14],
            # "Main_Q2c_SA_Impressive_Concept_D": [11],
        }

        lst_addin_MA_value = [
            # ['Q0a_RespondentID', 1004, 'Ma_SP', 1, 'Main_Q2c_SA_Impressive_Concept_A', 12],
            # ['Q0a_RespondentID', 1009, 'Ma_SP', 1, 'Main_Q2c_SA_Impressive_Concept_A', 11],
        ]

        # End Define structure------------------------------------------------------------------------------------------


        # CONVERT TO STACK----------------------------------------------------------------------------------------------
        df_data_stack, df_info_stack = self.convert_to_stack(df_data, df_info, id_col, sp_col, lst_scr, dict_main, lst_fc)
        df_data_stack_pack, df_info_stack_pack = self.convert_to_stack(df_data, df_info, id_col, pack_col, lst_scr, dict_pack, [])
        df_data_stack_price, df_info_stack_price = self.convert_to_stack(df_data, df_info, id_col, price_col, lst_scr, dict_price, [])

        # CONVERT TO STACK----------------------------------------------------------------------------------------------

        # OE RUNNING
        if dict_qre_OE_info:

            # ADD OE to Data stack--------------------------------------------------------------------------------------
            dict_qre_OE_product_info = dict()
            dict_qre_OE_pack_info = dict()
            dict_qre_OE_price_info = dict()

            for k, v in dict_qre_OE_info.items():
                if v[-1] == 'PRODUCT':
                    dict_qre_OE_product_info.update({k: v[:-1]})
                elif v[-1] == 'PACK':
                    dict_qre_OE_pack_info.update({k: v[:-1]})
                else:
                    dict_qre_OE_price_info.update({k: v[:-1]})

            df_data_stack[list(dict_qre_OE_product_info.keys())] = pd.DataFrame([[np.nan] * len(list(dict_qre_OE_product_info.keys()))], index=df_data_stack.index)
            df_data_stack_pack[list(dict_qre_OE_pack_info.keys())] = pd.DataFrame([[np.nan] * len(list(dict_qre_OE_pack_info.keys()))], index=df_data_stack_pack.index)
            df_data_stack_price[list(dict_qre_OE_price_info.keys())] = pd.DataFrame([[np.nan] * len(list(dict_qre_OE_price_info.keys()))], index=df_data_stack_price.index)

            # lst_OE_col = list(dict_qre_OE_info.keys())
            # df_data_stack[lst_OE_col] = pd.DataFrame([[np.nan] * len(lst_OE_col)], index=df_data_stack.index)

            # Remember edit this
            for item in lst_addin_OE_value:
                if item[2] == 'Ma_SP':
                    df_data_stack.loc[(df_data_stack[item[0]] == item[1]) & (df_data_stack[item[2]] == item[3]), [item[4]]] = [item[5]]
                elif item[2] == 'Ma_Pack':
                    df_data_stack_pack.loc[(df_data_stack_pack[item[0]] == item[1]) & (df_data_stack_pack[item[2]] == item[3]), [item[4]]] = [item[5]]
                else:
                    df_data_stack_price.loc[(df_data_stack_price[item[0]] == item[1]) & (df_data_stack_price[item[2]] == item[3]), [item[4]]] = [item[5]]

            # END ADD OE to Data stack----------------------------------------------------------------------------------

            # ADD OE to Info stack--------------------------------------------------------------------------------------
            df_info_stack = pd.concat([df_info_stack, pd.DataFrame(columns=['var_name', 'var_lbl', 'var_type', 'val_lbl'], data=list(dict_qre_OE_product_info.values()))], axis=0)
            df_info_stack_pack = pd.concat([df_info_stack_pack, pd.DataFrame(columns=['var_name', 'var_lbl', 'var_type', 'val_lbl'], data=list(dict_qre_OE_pack_info.values()))], axis=0)
            df_info_stack_price = pd.concat([df_info_stack_price, pd.DataFrame(columns=['var_name', 'var_lbl', 'var_type', 'val_lbl'], data=list(dict_qre_OE_price_info.values()))], axis=0)
            # END ADD OE to Info stack----------------------------------------------------------------------------------


        if lst_addin_MA_value:
            # ADD MA OE to Data stack-----------------------------------------------------------------------------------
            # Remember edit this
            for item in lst_addin_MA_value:

                # ['Q0a_RespondentID', 1004, 'Ma_SP', 1, 'Main_Q2c_SA_Impressive_Concept_A', 12],
                df_fil_info = df_info_stack.query(f"var_name.str.contains('{item[-2]}_[0-9]+')")
                df_fil_data = df_data_stack.query(f"{item[0]} == {item[1]} & {item[2]} == {item[3]}")

                for idx_row in df_fil_data.index:
                    for col_name in df_fil_info['var_name'].values.tolist():
                        if pd.isnull(df_fil_data.at[idx_row, col_name]):
                            df_data_stack.loc[idx_row, col_name] = item[-1]
                            break
                        else:
                            lst_clear_val = lst_clear_ma_oe_value[item[-2]]
                            aaa = df_fil_data.at[idx_row, col_name]

                            if df_fil_data.at[idx_row, col_name] in lst_clear_val:
                                df_data_stack.at[idx_row, col_name] = np.nan


            # END ADD MA OE to Data stack-------------------------------------------------------------------------------

        if dict_qre_net_info:
            # ADD MA NET CODE to df_info--------------------------------------------------------------------------------
            for key, val in dict_qre_net_info.items():
                df_info_stack.loc[df_info_stack['var_name'].str.contains(key), 'val_lbl'] = [val]
            # END ADD MA NET CODE to df_info----------------------------------------------------------------------------


        # REMEMBER RESET INDEX BEFORE RUN TABLES
        df_data_stack.reset_index(drop=True, inplace=True)
        df_info_stack.reset_index(drop=True, inplace=True)

        df_data_stack_pack.reset_index(drop=True, inplace=True)
        df_info_stack_pack.reset_index(drop=True, inplace=True)

        df_data_stack_price.reset_index(drop=True, inplace=True)
        df_info_stack_price.reset_index(drop=True, inplace=True)

        # PENALTY ANALYSIS----------------------------------------------------------------------------------------------
        # with pd.ExcelWriter('VN8327_ETIKA_df_data_stack.xlsx') as writer:
        #     df_data_stack.to_excel(writer, sheet_name='Data')
        #     df_info_stack.to_excel(writer, sheet_name='Map')

        df_temp = pd.DataFrame(
            columns=['Qre', 'Label', 'Ma_SP_Lbl', 'GroupCode', 'GroupCode_Pct', 'GroupCode_x_P1_Mean', 'JAR_x_P1_Mean', 'Penalty_Score', 'Pull_Down_Index'],
            data=[
                ['P4_JAR', 'AROMA', np.nan, 'B2B', np.nan, np.nan, np.nan, np.nan, np.nan],
                ['P4_JAR', 'AROMA', np.nan, 'T2B', np.nan, np.nan, np.nan, np.nan, np.nan],
                ['P5_JAR', 'FLAVOUR', np.nan, 'B2B', np.nan, np.nan, np.nan, np.nan, np.nan],
                ['P5_JAR', 'FLAVOUR', np.nan, 'T2B', np.nan, np.nan, np.nan, np.nan, np.nan],
                ['P6_JAR', 'SWEETNESS', np.nan, 'B2B', np.nan, np.nan, np.nan, np.nan, np.nan],
                ['P6_JAR', 'SWEETNESS', np.nan, 'T2B', np.nan, np.nan, np.nan, np.nan, np.nan],
                ['P7_JAR', 'SOURNESS', np.nan, 'B2B', np.nan, np.nan, np.nan, np.nan, np.nan],
                ['P7_JAR', 'SOURNESS', np.nan, 'T2B', np.nan, np.nan, np.nan, np.nan, np.nan],
                ['P8_JAR', 'TEXTURE', np.nan, 'B2B', np.nan, np.nan, np.nan, np.nan, np.nan],
                ['P8_JAR', 'TEXTURE', np.nan, 'T2B', np.nan, np.nan, np.nan, np.nan, np.nan],
            ])

        for item in ['Total', 'HCM', 'HN']:

            if item == 'Total':
                str_filter = 'S5 > 0'
            elif item == 'HCM':
                str_filter = 'S5 == 1'
            else:
                str_filter = 'S5 == 2'

            df_full_pen = pd.DataFrame()

            for k_sp, v_sp in {
                '101': 'CO',
                '102': 'CON',
                '103': 'YO',
                '104': 'PO',
                '105': 'CS',
                '106': 'CG',
                '107': 'PG',
                '108': 'PS',
                '109': 'GO',
                '110': 'GS',

                '102, 101, 105, 106': 'Calpis',  # CON, CO, CS, CG
                '109, 110': 'Goodday',  # GO, GS

            }.items():

                df_pen = df_temp.copy()

                df_pen['Ma_SP_Lbl'] = v_sp

                for qre in ['P4_JAR', 'P5_JAR', 'P6_JAR', 'P7_JAR', 'P8_JAR']:

                    # JAR
                    df_fil = df_data_stack.query(f"{qre}.isin([3]) & Ma_SP.isin([{k_sp}]) & {str_filter}").copy()
                    df_pen.loc[df_pen.eval(f"Qre == '{qre}'"), 'JAR_x_P1_Mean'] = df_fil['P1'].mean()
                    del df_fil

                    # B2B | T2B
                    for k_code, k_fil in {'B2B': '[1, 2]', 'T2B': '[4, 5]'}.items():
                        df_fil = df_data_stack.query(f"{qre}.isin({k_fil}) & Ma_SP.isin([{k_sp}]) & {str_filter}").copy()
                        df_pen.loc[df_pen.eval(f"Qre == '{qre}' & GroupCode == '{k_code}'"), 'GroupCode_x_P1_Mean'] = df_fil['P1'].mean()

                        df_fil_2 = df_data_stack.query(f"Ma_SP.isin([{k_sp}]) & {str_filter}").copy()
                        if k_code == 'B2B':
                            df_fil_2[qre].replace({1: 1, 2: 1, 3: 0, 4: 0, 5: 0}, inplace=True)
                        else:
                            df_fil_2[qre].replace({1: 0, 2: 0, 3: 0, 4: 1, 5: 1}, inplace=True)

                        # Pct
                        df_pen.loc[df_pen.eval(f"Qre == '{qre}' & GroupCode == '{k_code}'"), 'GroupCode_Pct'] = df_fil_2[qre].mean()

                        del df_fil
                        del df_fil_2


                df_pen['Penalty_Score'] = df_pen['JAR_x_P1_Mean'] - df_pen['GroupCode_x_P1_Mean']
                df_pen['Pull_Down_Index'] = df_pen['GroupCode_Pct'] * df_pen['Penalty_Score']

                df_full_pen = pd.concat([df_full_pen, df_pen], axis=0, ignore_index=True)






            with pd.ExcelWriter('VN8327_ETIKA_Penalty_Analysis_v1.xlsx', engine='openpyxl', mode='w' if item == 'Total' else 'a') as writer:
                df_full_pen.to_excel(writer, sheet_name=f'Penalty_Analysis_{item}')


        # SLD_LAYOUT_TITLE_AND_CONTENT = 1
        # prs = Presentation()
        # slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        # slide = prs.slides.add_slide(slide_layout)
        # prs.save('test.pptx')
        here = 1

        # PENALTY ANALYSIS----------------------------------------------------------------------------------------------


        # Export data tables--------------------------------------------------------------------------------------------
        if tables_format_file.filename:

            df_data_tbl, df_info_tbl = df_data_stack.copy(), df_info_stack.copy()

            df_data_tbl = pd.concat([df_data_tbl, df_data_stack_pack], axis=0)
            df_info_tbl = pd.concat([df_info_tbl, df_info_stack_pack], axis=0)

            df_data_tbl = pd.concat([df_data_tbl, df_data_stack_price], axis=0)
            df_info_tbl = pd.concat([df_info_tbl, df_info_stack_price], axis=0)

            df_data_tbl.sort_values(by=[id_col], inplace=True)
            df_data_tbl.reset_index(drop=True, inplace=True)

            df_info_tbl.drop_duplicates(subset=['var_name'], keep='first', inplace=True)
            df_info_tbl.reset_index(drop=True, inplace=True)

            # # ----------------------------------
            # df_data_tbl.to_excel('zzz_df_data_tbl.xlsx')
            # df_info_tbl.to_excel('zzz_df_info_tbl.xlsx')
            # # ----------------------------------

            # ADD MEAN & GROUP------------------------------------------------------------------------------------------
            self.logger.info('ADD MEAN & GROUP')

            lst_qre_mean = list()
            lst_qre_group = list()

            if dict_qre_group_mean:
                for key, val in dict_qre_group_mean.items():

                    if val['range']:
                        for i in val['range']:
                            lst_qre_mean.append([f'{key}_{i}', val['mean']])
                            lst_qre_group.append([f'{key}_{i}', val['group']])
                    else:
                        lst_qre_mean.append([key, val['mean']])
                        lst_qre_group.append([key, val['group']])

            # End ADD MEAN & GROUP--------------------------------------------------------------------------------------

            str_topline_file_name = self.str_file_name.replace('.xlsx', '_Topline.xlsx')
            DataTableGenerator.__init__(self, df_data=df_data_tbl, df_info=df_info_tbl,
                                        xlsx_name=str_topline_file_name, logger=self.logger,
                                        lst_qre_group=lst_qre_group, lst_qre_mean=lst_qre_mean, is_md=False)

            # lst_func_to_run = eval(tables_format_file.file.read())

            local_dic = locals()
            exec(tables_format_file.file.read(), globals(), local_dic)
            lst_func_to_run = local_dic['lst_func_to_run']


            self.run_tables_by_js_files(lst_func_to_run)
            self.format_sig_table()
        # End Export data tables----------------------------------------------------------------------------------------

        # Generate SAV files--------------------------------------------------------------------------------------------
        self.logger.info('Generate SAV files')

        df_info_stack = self.remove_net_code(df_info_stack)
        df_info_stack_pack = self.remove_net_code(df_info_stack_pack)
        df_info_stack_price = self.remove_net_code(df_info_stack_price)

        dict_dfs = {
            1: {
                'data': df_data_stack,
                'info': df_info_stack,
                'tail_name': 'product',
                'sheet_name': 'product',
                'is_recode_to_lbl': False,
            },
            2: {
                'data': df_data_stack_pack,
                'info': df_info_stack_pack,
                'tail_name': 'pack',
                'sheet_name': 'pack',
                'is_recode_to_lbl': False,
            },
            3: {
                'data': df_data_stack_price,
                'info': df_info_stack_price,
                'tail_name': 'price',
                'sheet_name': 'price',
                'is_recode_to_lbl': False,
            },
        }

        self.generate_multiple_sav_sps(dict_dfs=dict_dfs, is_md=False, is_export_xlsx=True)
        # END Generate SAV files----------------------------------------------------------------------------------------






    def remove_net_code(self, df_info: pd.DataFrame) -> pd.DataFrame:
        # Remove net_code to export sav---------------------------------------------------------------------------------
        df_info_without_net = df_info.copy()

        for idx in df_info_without_net.index:
            val_lbl = df_info_without_net.at[idx, 'val_lbl']

            if 'net_code' in val_lbl.keys():
                df_info_without_net.at[idx, 'val_lbl'] = self.unnetted_qre_val(val_lbl)

        return df_info_without_net
        # END Remove net_code to export sav-----------------------------------------------------------------------------