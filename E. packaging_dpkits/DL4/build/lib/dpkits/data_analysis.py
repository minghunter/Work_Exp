import pandas as pd
import numpy as np
import pingouin as pg





from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches





class DataAnalysis:
    def __init__(self, df_data: pd.DataFrame, df_info: pd.DataFrame):

        self.df_data = df_data
        self.df_info = df_info



    def penalty_analysis(self, dict_define_pen: dict, output_name: str):

        df_pen = pd.DataFrame(
            columns=['Section', 'Qre', 'Label', 'Ma_SP_Lbl', 'GroupCode', 'GroupCode_Pct', 'GroupCode_x_OL_Mean', 'JAR_x_OL_Mean', 'Penalty_Score', 'Pull_Down_Index'],
            data=[]
        )

        df_info = self.df_info.copy()

        for k_sec, v_sec in dict_define_pen.items():
            print(f'Processing penalty analysis - {k_sec}')

            df_data = self.df_data.query(v_sec.get('query')).copy() if v_sec.get('query') else self.df_data.copy()

            for k_sp, v_sp in df_info.loc[df_info.eval(f"var_name == '{v_sec['prod_pre']}'"), 'val_lbl'].values[0].items():

                df_fil = df_data.query(f"{v_sec['prod_pre']}.isin([{k_sp}])")

                for k_jar, v_jar in v_sec['jar_qres'].items():

                    jar_ol_mean = df_fil.loc[df_fil.eval(f"{k_jar}.isin({v_jar['jar']['code']})"), v_sec['ol_qre']].mean()

                    for grp in ['b2b', 't2b']:
                        grp_count = df_fil.loc[df_fil.eval(f"{k_jar}.isin({v_jar[grp]['code']})"), k_jar].count()
                        grp_base = df_fil.loc[df_fil.eval(f"{k_jar}.notnull()"), k_jar].count()

                        if not grp_base:
                            continue


                        grp_pct = grp_count / grp_base
                        grp_ol_mean = df_fil.loc[df_fil.eval(f"{k_jar}.isin({v_jar[grp]['code']})"), v_sec['ol_qre']].mean()
                        pen_score = jar_ol_mean - grp_ol_mean

                        dict_pen_data_row = {
                            'Section': k_sec,
                            'Qre': k_jar,
                            'Label': v_jar['label'],
                            'Ma_SP_Lbl': v_sp,
                            'GroupCode': v_jar[grp]['label'],
                            'GroupCode_Pct': grp_pct,
                            'GroupCode_x_OL_Mean': grp_ol_mean,
                            'JAR_x_OL_Mean': jar_ol_mean,
                            'Penalty_Score': pen_score,
                            'Pull_Down_Index': grp_pct * pen_score,
                        }

                        if df_pen.empty:
                            df_pen = pd.DataFrame(columns=list(dict_pen_data_row.keys()), data=[dict_pen_data_row.values()])
                        else:
                            df_pen = pd.concat([df_pen, pd.DataFrame(columns=list(dict_pen_data_row.keys()), data=[dict_pen_data_row.values()])], axis=0, ignore_index=True)

        with pd.ExcelWriter(f'{output_name}.xlsx', engine='openpyxl') as writer:
            df_pen.to_excel(writer, sheet_name=f'Penalty_Analysis')



        # HERE
        # # --------------------------------------------------------------------------------------------------------------
        # # create presentation with 1 slide ------
        # prs = Presentation()
        # slide = prs.slides.add_slide(prs.slide_layouts[5])
        #
        # # Export ppt chart after penalty score table
        # chart_data = XyChartData()
        #
        # series_1 = chart_data.add_series('Model 1')
        # series_1.add_data_point(0.7, 2.7)
        #
        # series_2 = chart_data.add_series('Model 2')
        # series_2.add_data_point(1.3, 3.7)
        #
        # series_2 = chart_data.add_series('Model 3')
        # series_2.add_data_point(0.6, 1.3)
        #
        # x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(4.5)
        #
        # chart = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data)
        #
        # prs.save('chart-01-test.pptx')









    def linear_regression(self, dict_define_linear: dict, output_name: str):

        # Single: y = b + a*x
        # Multiple: y = b + a1*x1 + a2*x2 + ... + an*xn

        with pd.ExcelWriter(f'{output_name}.xlsx', engine='openpyxl') as writer:
            for k_ln, v_ln in dict_define_linear.items():
                print(f'Processing linear regression - {k_ln}')
                df_data = self.df_data.query(str_query).copy() if v_ln['str_query'] else self.df_data.copy()

                df_data.loc[:, 'dep_var'] = df_data.loc[:, v_ln['dependent_vars']].mean(axis=1)

                ln_reg = pg.linear_regression(df_data.loc[:, v_ln['explanatory_vars']], df_data['dep_var'])

                ln_reg.to_excel(writer, sheet_name=k_ln)



    def logistic_regression(self, dict_define_logistic: dict, output_name: str):
        
        # add calculate function here
        



        pass




    # MORE ANALYSIS HERE

